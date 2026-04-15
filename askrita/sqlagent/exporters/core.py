# Copyright 2026 CVS Health and/or one of its affiliates
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
#
# This file uses the following unmodified third-party packages,
# each retaining its original copyright and license:
#   python-pptx (MIT)
#   reportlab (BSD-3-Clause)

"""Core export functionality for PPTX and PDF generation."""

import io
import logging
from datetime import datetime

# Optional export dependencies
try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        Image,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

from askrita.sqlagent.State import WorkflowState

from ...exceptions import ExportError
from .chart_generator import (
    add_native_pptx_chart,
    generate_chart_bytes,
    get_chart_data_for_export,
)
from .models import ExportSettings


def _header_for_key(chart_data, i: int, key: str) -> str:
    """Return a chart-metadata-enhanced header for column index ``i``."""
    if i == 0 and hasattr(chart_data, "xAxisLabel") and chart_data.xAxisLabel:
        return chart_data.xAxisLabel
    if (
        i == 1
        and hasattr(chart_data, "datasets")
        and chart_data.datasets
        and hasattr(chart_data.datasets[0], "label")
        and chart_data.datasets[0].label
    ):
        return chart_data.datasets[0].label
    if hasattr(chart_data, "yAxisLabel") and chart_data.yAxisLabel and i > 0:
        return chart_data.yAxisLabel
    return key.title()


def _dataset_label(dataset, chart_data) -> str:
    """Return the best label for a single dataset entry."""
    if hasattr(dataset, "label") and dataset.label:
        return dataset.label
    if hasattr(chart_data, "yAxisLabel") and chart_data.yAxisLabel:
        return chart_data.yAxisLabel
    return "Value"


def _first_column_header(chart_data, fallback_headers) -> str:
    """Return the header for the first (X-axis) column."""
    if hasattr(chart_data, "xAxisLabel") and chart_data.xAxisLabel:
        return chart_data.xAxisLabel
    if fallback_headers:
        return fallback_headers[0]
    return "Category"


def _dataset_headers(chart_data) -> list:
    """Return the list of value-column headers derived from chart datasets or y-axis label."""
    if hasattr(chart_data, "datasets") and chart_data.datasets:
        return [_dataset_label(ds, chart_data) for ds in chart_data.datasets]
    if hasattr(chart_data, "yAxisLabel") and chart_data.yAxisLabel:
        return [chart_data.yAxisLabel]
    return []


def _generate_table_headers_from_chart_data(chart_data, results, fallback_headers):
    """
    Generate intelligent table headers using chart metadata (axis labels, dataset labels).

    Args:
        chart_data: UniversalChartData with rich metadata
        results: Raw data results
        fallback_headers: Fallback headers if chart metadata isn't available

    Returns:
        List of meaningful column headers
    """
    if not chart_data or not results:
        return fallback_headers

    headers = [_first_column_header(chart_data, fallback_headers)] + _dataset_headers(
        chart_data
    )

    if results and isinstance(results[0], dict):
        result_keys = list(results[0].keys())
        if len(headers) != len(result_keys):
            return [
                _header_for_key(chart_data, i, key) for i, key in enumerate(result_keys)
            ]

    return headers if headers else fallback_headers


logger = logging.getLogger(__name__)


def _pptx_add_title_slide(prs, settings, primary_color, secondary_color, dark_gray):
    """Add the title slide (slide 1) to the presentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = settings.title
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = secondary_color
    title_para.alignment = PP_ALIGN.CENTER

    subtitle = slide.placeholders[1]
    subtitle.text = f"{settings.company_name}\n\nData Analytics Report\nGenerated {datetime.now().strftime('%B %d, %Y')}"
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = dark_gray
        paragraph.alignment = PP_ALIGN.CENTER

    footer_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7), Inches(10), Inches(0.5)
    )
    footer_shape.fill.solid()
    footer_shape.fill.fore_color.rgb = primary_color
    footer_shape.line.fill.background()


def _pptx_add_summary_slide(prs, output_state, primary_color, dark_gray):
    """Add the Executive Summary slide (slide 2) and return the slide layout used."""
    content_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_slide_layout)

    title = slide.shapes.title
    title.text = "Executive Summary"
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color

    content_placeholder = slide.placeholders[1]
    slide.shapes._spTree.remove(content_placeholder._element)

    left_column = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5)
    )
    left_frame = left_column.text_frame
    left_frame.word_wrap = True
    p1 = left_frame.paragraphs[0]
    p1.text = "Question & Answer"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = primary_color
    p1.space_after = Pt(16)

    if output_state.question:
        p2 = left_frame.add_paragraph()
        p2.text = "Question:"
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = dark_gray
        p2.space_after = Pt(6)
        p3 = left_frame.add_paragraph()
        p3.text = output_state.question
        p3.font.size = Pt(13)
        p3.font.color.rgb = dark_gray
        p3.space_after = Pt(16)

    if output_state.answer:
        p4 = left_frame.add_paragraph()
        p4.text = "Answer:"
        p4.font.size = Pt(14)
        p4.font.bold = True
        p4.font.color.rgb = dark_gray
        p4.space_after = Pt(6)
        p5 = left_frame.add_paragraph()
        p5.text = output_state.answer
        p5.font.size = Pt(13)
        p5.font.color.rgb = dark_gray

    right_column = slide.shapes.add_textbox(
        Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5)
    )
    right_frame = right_column.text_frame
    right_frame.word_wrap = True
    p1_right = right_frame.paragraphs[0]
    p1_right.text = "Technical Details"
    p1_right.font.size = Pt(18)
    p1_right.font.bold = True
    p1_right.font.color.rgb = primary_color
    p1_right.space_after = Pt(16)

    if output_state.sql_reason:
        p2_right = right_frame.add_paragraph()
        p2_right.text = "SQL Generation:"
        p2_right.font.size = Pt(14)
        p2_right.font.bold = True
        p2_right.font.color.rgb = dark_gray
        p2_right.space_after = Pt(6)
        p3_right = right_frame.add_paragraph()
        p3_right.text = output_state.sql_reason
        p3_right.font.size = Pt(12)
        p3_right.font.color.rgb = dark_gray
        p3_right.space_after = Pt(12)

    if output_state.visualization_reason:
        p4_right = right_frame.add_paragraph()
        p4_right.text = "Visualization Choice:"
        p4_right.font.size = Pt(14)
        p4_right.font.bold = True
        p4_right.font.color.rgb = dark_gray
        p4_right.space_after = Pt(6)
        p5_right = right_frame.add_paragraph()
        p5_right.text = output_state.visualization_reason
        p5_right.font.size = Pt(12)
        p5_right.font.color.rgb = dark_gray

    return content_slide_layout


def _pptx_add_viz_slide(prs, output_state, primary_color, dark_gray, light_gray):
    """Add the Visualization slide (slide 3) and return (chart_data, chart_type)."""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = light_gray
    title_bg.line.color.rgb = primary_color
    title_bg.line.width = Pt(2)

    chart_data, chart_type = get_chart_data_for_export(output_state)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    p1 = title_frame.paragraphs[0]
    if hasattr(chart_data, "title") and chart_data.title:
        p1.text = chart_data.title
    else:
        display_chart_type = (
            chart_type if chart_type and chart_type != "none" else "Data"
        )
        p1.text = f"Data Visualization: {display_chart_type.title()} Chart"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = primary_color
    p1.alignment = PP_ALIGN.LEFT

    if output_state.question:
        p2 = title_frame.add_paragraph()
        p2.text = f"Question: {output_state.question}"
        p2.font.size = Pt(16)
        p2.font.color.rgb = dark_gray
        p2.alignment = PP_ALIGN.LEFT

    try:
        chart_added = add_native_pptx_chart(
            slide,
            chart_data,
            chart_type,
            Inches(0.5),
            Inches(1.8),
            Inches(9),
            Inches(5.2),
        )
        if chart_added:
            logger.info("Native PowerPoint chart added successfully")
        else:
            logger.error("Native chart generation failed")
    except Exception as chart_error:
        logger.error(f"Chart generation error: {chart_error}")

    if hasattr(chart_data, "yAxes") and chart_data.yAxes and len(chart_data.yAxes) > 1:
        note_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.2), Inches(9), Inches(0.5)
        )
        note_frame = note_box.text_frame
        note_frame.word_wrap = True
        note_frame.clear()
        note_para = note_frame.paragraphs[0]
        note_para.text = (
            "Multi-Axis Chart Detected: To add secondary Y-axis, right-click series "
            "-> Format Data Series -> Plot Series On -> Secondary Axis"
        )
        note_para.font.size = Pt(10)
        note_para.font.italic = True
        note_para.font.color.rgb = RGBColor(100, 100, 100)
        note_para.alignment = PP_ALIGN.LEFT
        logger.info("Added multi-axis configuration note to slide")

    return blank_slide_layout, chart_data, chart_type


def _pptx_fallback_headers_from_row(first_row, sql_query):
    """Derive fallback column headers from a result row and optional SQL query."""
    if isinstance(first_row, dict):
        return list(first_row.keys())
    if isinstance(first_row, (tuple, list)):
        fallback_headers = [f"Column_{i + 1}" for i in range(len(first_row))]
        if sql_query:
            import re

            select_match = re.search(
                r"SELECT\s+(.*?)\s+FROM", sql_query, re.IGNORECASE | re.DOTALL
            )
            if select_match:
                columns = [
                    col.strip().split(" AS ")[-1].split(".")[-1].strip('`"[]')
                    for col in select_match.group(1).split(",")
                ]
                if len(columns) == len(first_row):
                    logger.info(
                        f"Extracted {len(columns)} column names from SQL for PPTX table"
                    )
                    return columns
        return fallback_headers
    return []


def _pptx_row_values(row_data, headers, fallback_headers):
    """Extract ordered cell values from a result row using fallback header keys."""
    if isinstance(row_data, dict):
        values = []
        for col_idx, intelligent_header in enumerate(headers):
            if col_idx < len(fallback_headers):
                values.append(row_data.get(fallback_headers[col_idx], ""))
            else:
                values.append(row_data.get(intelligent_header, ""))
        return values
    if isinstance(row_data, (tuple, list)):
        return list(row_data)
    return [str(row_data)]


def _pptx_add_data_table_slide(
    prs, output_state, content_slide_layout, primary_color, dark_gray, light_gray
):
    """Add the Data Table slide (slide 4) if data is available."""
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    title.text = "Detailed Data"
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    content_placeholder = slide.placeholders[1]
    slide.shapes._spTree.remove(content_placeholder._element)

    table_data = output_state.results[:15]
    if not table_data:
        return

    fallback_headers = _pptx_fallback_headers_from_row(
        table_data[0], output_state.sql_query
    )
    chart_data, _ = get_chart_data_for_export(output_state)
    if chart_data:
        headers = _generate_table_headers_from_chart_data(
            chart_data, table_data, fallback_headers
        )
        logger.info(f"Generated PPTX headers from chart metadata: {headers}")
    else:
        headers = fallback_headers
        logger.info(f"Using fallback headers for PPTX: {headers}")

    if not headers:
        return

    table_shape = slide.shapes.add_table(
        len(table_data) + 1,
        len(headers),
        Inches(0.5),
        Inches(1.5),
        Inches(9),
        Inches(5.5),
    )
    table = table_shape.table

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = primary_color
        para = cell.text_frame.paragraphs[0]
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.font.bold = True
        para.font.size = Pt(12)
        para.alignment = PP_ALIGN.CENTER

    for row_idx, row_data in enumerate(table_data, 1):
        values = _pptx_row_values(row_data, headers, fallback_headers)
        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value) if value is not None else ""
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = light_gray
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.font.color.rgb = dark_gray
            para.alignment = PP_ALIGN.CENTER


def _pptx_add_sql_slide(
    prs, output_state, content_slide_layout, primary_color, dark_gray, light_gray
):
    """Add the SQL Query slide (slide 5)."""
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    title.text = "SQL Query"
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    content_placeholder = slide.placeholders[1]
    slide.shapes._spTree.remove(content_placeholder._element)

    sql_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    sql_frame = sql_box.text_frame
    sql_frame.word_wrap = True
    sql_frame.text = output_state.sql_query
    sql_para = sql_frame.paragraphs[0]
    sql_para.font.name = "Consolas"
    sql_para.font.size = Pt(10)
    sql_para.font.color.rgb = dark_gray
    sql_box.fill.solid()
    sql_box.fill.fore_color.rgb = light_gray
    sql_box.line.color.rgb = dark_gray
    sql_box.line.width = Pt(1)


def _pptx_add_followup_slide(
    prs, output_state, content_slide_layout, primary_color, dark_gray
):
    """Add the Follow-up Questions slide (slide 6)."""
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    title.text = "Further Analysis"
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color

    content = slide.placeholders[1]
    content.text = ""
    text_frame = content.text_frame
    text_frame.clear()
    p1 = text_frame.paragraphs[0]
    p1.text = "Consider exploring these follow-up analyses:"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = dark_gray
    p1.space_after = Pt(16)

    for question in output_state.followup_questions:
        p = text_frame.add_paragraph()
        p.text = f" {question}"
        p.level = 1
        p.font.size = Pt(16)
        p.font.color.rgb = dark_gray
        p.space_after = Pt(12)


def _pptx_add_closing_slide(prs, settings, primary_color, light_gray):
    """Add the closing Thank You slide."""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(2), Inches(10), Inches(4)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = light_gray
    bg_shape.line.color.rgb = primary_color
    bg_shape.line.width = Pt(2)

    thank_you_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = f"Thank you\n\nReport generated by {settings.company_name}\nAskRITA Analytics Platform"
    for paragraph in thank_you_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = primary_color
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER


def create_pptx_export(output_state: WorkflowState, settings: ExportSettings) -> bytes:
    """
    Create a professional PowerPoint presentation from query results.

    Args:
        output_state: Complete OutputState from workflow.query() or workflow.chat()
        settings: Export customization settings

    Returns:
        bytes: PPTX file bytes

    Raises:
        ExportError: If PPTX generation fails
        ImportError: If python-pptx is not installed
    """
    if not PPTX_AVAILABLE:
        raise ImportError(
            "PPTX export requires python-pptx. Install with: pip install askrita[exports]"
        )

    try:
        prs = Presentation()

        primary_color = RGBColor(*settings.brand_primary_color)
        secondary_color = RGBColor(*settings.brand_secondary_color)
        dark_gray = RGBColor(64, 64, 64)
        light_gray = RGBColor(242, 242, 242)

        _pptx_add_title_slide(prs, settings, primary_color, secondary_color, dark_gray)
        content_slide_layout = _pptx_add_summary_slide(
            prs, output_state, primary_color, dark_gray
        )
        _pptx_add_viz_slide(prs, output_state, primary_color, dark_gray, light_gray)

        if settings.include_data_table and output_state.results:
            _pptx_add_data_table_slide(
                prs,
                output_state,
                content_slide_layout,
                primary_color,
                dark_gray,
                light_gray,
            )

        if settings.include_sql and output_state.sql_query:
            _pptx_add_sql_slide(
                prs,
                output_state,
                content_slide_layout,
                primary_color,
                dark_gray,
                light_gray,
            )

        if output_state.followup_questions:
            _pptx_add_followup_slide(
                prs, output_state, content_slide_layout, primary_color, dark_gray
            )

        _pptx_add_closing_slide(prs, settings, primary_color, light_gray)

        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        logger.info("PPTX export completed successfully")
        return pptx_buffer.getvalue()

    except Exception as e:
        logger.error(f"PPTX generation failed: {e}")
        raise ExportError(f"PPTX generation failed: {str(e)}")


def _pdf_append_chart(story, output_state, settings):
    """Generate a chart image and append it to the PDF story list."""
    chart_data, chart_type = get_chart_data_for_export(output_state)
    if not (chart_data and chart_type and chart_type != "none"):
        return
    try:
        chart_title = (
            chart_data.title
            if hasattr(chart_data, "title") and chart_data.title
            else "Data Analysis Results"
        )
        chart_bytes = generate_chart_bytes(
            chart_data, chart_type, chart_title, settings.chart_style
        )
        if chart_bytes:
            styles = getSampleStyleSheet()
            story.append(Paragraph("Data Visualization", styles["Heading2"]))
            try:
                story.append(
                    Image(io.BytesIO(chart_bytes), width=7 * inch, height=4.5 * inch)
                )
                story.append(Spacer(1, 20))
                logger.info("Chart added to PDF")
            except Exception as img_error:
                logger.error(f"Failed to add chart to PDF: {img_error}")
    except Exception as chart_error:
        logger.error(f"Chart generation error: {chart_error}")


def _pdf_append_data_table(story, output_state):
    """Build and append a data table to the PDF story list."""
    styles = getSampleStyleSheet()
    story.append(Paragraph("Data Table", styles["Heading2"]))

    table_data = output_state.results[:20]
    if len(output_state.results) > 20:
        story.append(
            Paragraph(
                f"Showing first 20 rows of {len(output_state.results)} total results",
                styles["Normal"],
            )
        )

    if not table_data:
        return

    first_row = table_data[0]
    if not isinstance(first_row, dict):
        return

    fallback_headers = list(first_row.keys())
    chart_data, _ = get_chart_data_for_export(output_state)
    if chart_data:
        headers = _generate_table_headers_from_chart_data(
            chart_data, table_data, fallback_headers
        )
        logger.info(f"Generated PDF headers from chart metadata: {headers}")
    else:
        headers = fallback_headers
        logger.info(f"Using fallback headers for PDF: {headers}")

    table_rows = [[str(h) for h in headers]]
    for row_dict in table_data:
        table_rows.append(
            [
                str(
                    row_dict.get(
                        (
                            fallback_headers[i]
                            if i < len(fallback_headers)
                            else headers[i]
                        ),
                        "",
                    )
                )
                for i in range(len(headers))
            ]
        )

    table = Table(table_rows)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, 0), 12),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
                ("BACKGROUND", (0, 1), (-1, -1), colors.beige),
                ("GRID", (0, 0), (-1, -1), 1, colors.black),
            ]
        )
    )
    story.append(table)
    story.append(Spacer(1, 20))


def create_pdf_export(output_state: WorkflowState, settings: ExportSettings) -> bytes:
    """
    Create PDF report from query results.

    Args:
        output_state: Complete OutputState from workflow.query() or workflow.chat()
        settings: Export customization settings

    Returns:
        bytes: PDF file bytes

    Raises:
        ExportError: If PDF generation fails
        ImportError: If reportlab is not installed
    """
    if not PDF_AVAILABLE:
        raise ImportError(
            "PDF export requires reportlab. Install with: pip install askrita[exports]"
        )

    try:
        pdf_buffer = io.BytesIO()
        doc = SimpleDocTemplate(pdf_buffer, pagesize=A4, topMargin=1 * inch)
        styles = getSampleStyleSheet()
        story = []

        title_style = ParagraphStyle(
            "CustomTitle",
            parent=styles["Heading1"],
            fontSize=24,
            alignment=TA_CENTER,
            spaceAfter=30,
        )
        story.append(Paragraph(settings.title, title_style))
        story.append(Paragraph(f"{settings.company_name}", styles["Normal"]))
        story.append(
            Paragraph(
                f"Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}",
                styles["Normal"],
            )
        )
        story.append(Spacer(1, 20))

        story.append(Paragraph("Query Analysis", styles["Heading2"]))
        if output_state.question:
            story.append(
                Paragraph(f"<b>Question:</b> {output_state.question}", styles["Normal"])
            )
            story.append(Spacer(1, 6))
        if output_state.answer:
            story.append(
                Paragraph(f"<b>Answer:</b> {output_state.answer}", styles["Normal"])
            )
        if output_state.sql_reason:
            story.append(
                Paragraph(
                    f"<b>SQL Generation Reasoning:</b> {output_state.sql_reason}",
                    styles["Normal"],
                )
            )
        if output_state.visualization_reason:
            story.append(
                Paragraph(
                    f"<b>Visualization Choice:</b> {output_state.visualization_reason}",
                    styles["Normal"],
                )
            )
        story.append(Spacer(1, 12))

        if settings.include_sql and output_state.sql_query:
            story.append(Paragraph("<b>SQL Query:</b>", styles["Normal"]))
            sql_style = ParagraphStyle(
                "SQL",
                parent=styles["Code"],
                fontSize=10,
                leftIndent=20,
                backgroundColor=colors.lightgrey,
            )
            story.append(Paragraph(output_state.sql_query, sql_style))
            story.append(Spacer(1, 20))

        _pdf_append_chart(story, output_state, settings)

        if settings.include_data_table and output_state.results:
            _pdf_append_data_table(story, output_state)

        if output_state.followup_questions:
            story.append(Paragraph("Follow-up Questions", styles["Heading2"]))
            for i, question in enumerate(output_state.followup_questions, 1):
                story.append(Paragraph(f"{i}. {question}", styles["Normal"]))
            story.append(Spacer(1, 12))

        doc.build(story)
        pdf_buffer.seek(0)

        logger.info("PDF export completed successfully")
        return pdf_buffer.getvalue()

    except Exception as e:
        logger.error(f"PDF generation failed: {e}")
        raise ExportError(f"PDF generation failed: {str(e)}")
